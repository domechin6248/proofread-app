import streamlit as st
import pandas as pd
import docx
from docx.shared import RGBColor, Pt
from docx.oxml.ns import qn
import openpyxl
from openpyxl.styles import Font
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
import os
import re
from io import BytesIO
import pdfplumber

# ページ設定
st.set_page_config(page_title="川内JC 統一ルール修正ツール", page_icon="⚓", layout="wide")
st.title("⚓ 2026年度 川内JC 統一ルール 修正システム")

# 1. 色の設定
color_option = st.selectbox("修正箇所の文字色（Word/Excel/PPT用）", ["赤", "青", "緑", "黒"], index=0)
color_map = {"赤": (255, 0, 0), "青": (0, 0, 255), "緑": (0, 128, 0), "黒": (0, 0, 0)}
selected_rgb = color_map[color_option]

# 2. ルールの読み込み
@st.cache_data(ttl=1)
def load_rules():
    if os.path.exists('rules.csv'):
        encodings = ['utf-8-sig', 'shift-jis', 'utf-8', 'cp932']
        for enc in encodings:
            try:
                df = pd.read_csv('rules.csv', encoding=enc)
                if '類義語' in df.columns and '統一語句' in df.columns:
                    df = df.dropna(subset=['類義語', '統一語句'])
                    df['類義語'] = df['類義語'].astype(str).str.strip()
                    df['統一語句'] = df['統一語句'].astype(str).str.strip()
                    df['len'] = df['類義語'].str.len()
                    df = df.sort_values('len', ascending=False)
                    return dict(zip(df['類義語'], df['統一語句']))
            except:
                continue
    return {}

rules_dict = load_rules()

# 3. 修正・熟語保護・英数字半角化ロジック
def apply_rules_to_text(target_text, rules, for_reporting=False):
    keep_words = [
        "会員に成長する機会", "会員拡大運動", "会員拡大", "正会員", "新入会員",
        "日本の青年会議所は", "希望をもたらす変革の起点として", 
        "輝く個性が調和する未来を描き", "社会の課題を解決することで", 
        "持続可能な地域を創ることを誓う", "われわれ JAYCEE は", "われわれJAYCEEは",
        "志高き組織ビジョン", "志高き人材育成ビジョン", "志高きまち創造ビジョン"
    ]
    
    for k, v in rules.items():
        if k == v and str(k) not in keep_words:
            keep_words.append(str(k))
            
    keep_words = sorted(keep_words, key=len, reverse=True)
    protected_text = target_text
    placeholders = {}
    
    for i, word in enumerate(keep_words):
        if word in protected_text:
            placeholder = f"《《保{i:04d}護》》"
            placeholders[placeholder] = word
            protected_text = protected_text.replace(word, placeholder)

    segments = [(protected_text, protected_text, False)]
    for wrong, right in rules.items():
        if wrong == right or not wrong: continue
        new_segments = []
        for orig, curr, already_fixed in segments:
            if already_fixed or str(wrong) not in curr:
                new_segments.append((orig, curr, already_fixed))
                continue
            parts = curr.split(str(wrong))
            for j, part in enumerate(parts):
                if part != "":
                    new_segments.append((part, part, False))
                if j < len(parts) - 1:
                    new_segments.append((str(wrong), str(right), True))
        segments = new_segments

    restored_segments = []
    for orig, curr, is_fixed in segments:
        temp_orig = orig
        temp_curr = curr
        if not is_fixed:
            for placeholder, original_word in placeholders.items():
                temp_orig = temp_orig.replace(placeholder, original_word)
                temp_curr = temp_curr.replace(placeholder, original_word)
        restored_segments.append((temp_orig, temp_curr, is_fixed))

    ZEN_ALNUM = "ＡＢＣＤＥＦＧＨＩＪＫＬＭＮＯＰＱＲＳＴＵＶＷＸＹＺａｂｃｄｅｆｇｈｉｊｋｌｍｎｏｐｑｒｓｔｕｖｗｘｙｚ０１２３４５６７８９"
    HAN_ALNUM = "ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789"
    ZEN2HAN_MAP = str.maketrans(ZEN_ALNUM, HAN_ALNUM)
    alnum_pattern = r'([A-Za-z0-9Ａ-Ｚａ-ｚ０-９]+)'
    final_segments = []
    
    for orig, curr, is_fixed in restored_segments:
        if for_reporting and is_fixed:
            def replacer(match):
                return match.group(1).translate(ZEN2HAN_MAP)
            new_curr = re.sub(alnum_pattern, replacer, curr)
            has_alnum = bool(re.search(alnum_pattern, new_curr))
            final_segments.append((orig, new_curr, True, has_alnum))
        else:
            parts = re.split(alnum_pattern, curr)
            for i, part in enumerate(parts):
                if not part: continue
                if i % 2 == 1:
                    half_part = part.translate(ZEN2HAN_MAP)
                    was_converted = (half_part != part)
                    if for_reporting:
                        final_segments.append((part, half_part, was_converted, True))
                    else:
                        final_segments.append((orig, half_part, is_fixed or was_converted, True))
                else:
                    if for_reporting:
                        final_segments.append((part, part, False, False))
                    else:
                        final_segments.append((orig, part, is_fixed, False))
                        
    return final_segments

# --- 網掛け（背景色）検知機能 ---
def is_word_shaded(para):
    try:
        pPr = para._p.pPr
        if pPr is not None:
            shd = pPr.find(qn('w:shd'))
            if shd is not None:
                val = shd.get(qn('w:val'))
                fill = shd.get(qn('w:fill'))
                if val and val != 'clear': return True
                if fill and fill not in ['auto', 'FFFFFF', 'clear']: return True
                
        parent = para._p.getparent()
        if parent is not None and parent.tag.endswith('tc'):
            tcPr = parent.find(qn('w:tcPr'))
            if tcPr is not None:
                shd = tcPr.find(qn('w:shd'))
                if shd is not None:
                    val = shd.get(qn('w:val'))
                    fill = shd.get(qn('w:fill'))
                    if val and val != 'clear': return True
                    if fill and fill not in ['auto', 'FFFFFF', 'clear']: return True
    except:
        pass
    return False

# --- ファイル修正用関数 ---
def repair_docx(file, rules, rgb):
    doc = docx.Document(file)
    
    def process_paragraphs(paragraphs):
        for para in paragraphs:
            is_shaded = is_word_shaded(para)
            orig_bold, orig_size = None, None
            if para.runs and para.runs[0].font:
                orig_bold = para.runs[0].font.bold
                orig_size = para.runs[0].font.size

            # 【新機能】段落内の「リンク」と「普通の文字」を分けて抽出
            elements = []
            current_text = ""
            
            try:
                # WordのXMLタグを直接走査し、リンクをカプセルごと避難させる
                for child in list(para._p):
                    if child.tag.endswith('hyperlink'):
                        if current_text:
                            elements.append({"type": "text", "content": current_text})
                            current_text = ""
                        # リンクをXMLから引き剥がして保存
                        elements.append({"type": "link", "element": child})
                        para._p.remove(child)
                    elif child.tag.endswith('r') or child.tag.endswith('ins'):
                        text = "".join([t.text for t in child.xpath('.//w:t') if t.text])
                        current_text += text
                        para._p.remove(child)
                
                if current_text:
                    elements.append({"type": "text", "content": current_text})
            except:
                elements = [{"type": "text", "content": para.text}]
                para.text = ""

            # 普通の文字にはルールを適用し、リンクはそのまま再配置する
            for el in elements:
                if el["type"] == "text":
                    if not el["content"]: continue
                    parts = apply_rules_to_text(el["content"], rules, for_reporting=False)
                    for orig, curr, is_fixed, is_alnum in parts:
                        run = para.add_run(curr)
                        run.font.name = 'ＭＳ 明朝'
                        run._element.rPr.rFonts.set(qn('w:eastAsia'), 'ＭＳ 明朝')
                        
                        if is_shaded:
                            if orig_size is not None: run.font.size = orig_size
                            if orig_bold is not None: run.font.bold = orig_bold
                        else:
                            run.font.size = Pt(10.5)
                            
                        if is_fixed:
                            run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                            run.bold = False
                            
                elif el["type"] == "link":
                    # 避難させていたリンクのカプセルを段落に戻す（機能・見た目すべて完全維持）
                    para._p.append(el["element"])

    process_paragraphs(doc.paragraphs)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                process_paragraphs(cell.paragraphs)

    out_io = BytesIO()
    doc.save(out_io)
    return out_io.getvalue()

def repair_xlsx(file, rules, rgb):
    wb = openpyxl.load_workbook(file)
    hex_color = f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    # 【新機能】セルにリンクが設定されている場合は、一切修正せず完全スルーする
                    if cell.hyperlink:
                        continue
                        
                    is_shaded = False
                    if cell.fill and cell.fill.patternType and cell.fill.patternType != 'none':
                        if cell.fill.fgColor.rgb and cell.fill.fgColor.rgb not in ['00000000', 'FFFFFFFF', '00FFFFFF']:
                            is_shaded = True

                    orig_bold = cell.font.bold if cell.font else False
                    orig_size = cell.font.size if cell.font else 11

                    parts = apply_rules_to_text(cell.value, rules, for_reporting=False)
                    
                    if any(p[2] for p in parts) or any(p[3] for p in parts) or not is_shaded:
                        cell.value = "".join([p[1] for p in parts])
                        is_fixed_present = any(p[2] for p in parts)
                        new_color = hex_color if is_fixed_present else (cell.font.color if cell.font else None)
                        new_size = orig_size if is_shaded else 10.5
                        new_bold = orig_bold if is_shaded else False
                        
                        cell.font = Font(name='ＭＳ 明朝', size=new_size, color=new_color, bold=new_bold)
                        
    out_io = BytesIO()
    wb.save(out_io)
    return out_io.getvalue()

def repair_pptx(file, rules, rgb):
    prs = Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            is_shaded = False
            if hasattr(shape, "fill") and shape.fill.type == 1:
                is_shaded = True
                
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # 【新機能】段落内にリンクが含まれる場合はスルーして保護する
                    has_link = any(hasattr(run, "hyperlink") and run.hyperlink and run.hyperlink.address for run in paragraph.runs)
                    if has_link:
                        continue
                        
                    combined_text = "".join(run.text for run in paragraph.runs)
                    parts = apply_rules_to_text(combined_text, rules, for_reporting=False)
                    
                    if any(p[2] for p in parts) or any(p[3] for p in parts) or not is_shaded:
                        paragraph.text = ""
                        for orig, curr, is_fixed, is_alnum in parts:
                            new_run = paragraph.add_run()
                            new_run.text = curr
                            new_run.font.name = 'ＭＳ 明朝'
                            if not is_shaded:
                                new_run.font.size = Pt(10.5)
                            if is_fixed:
                                new_run.font.color.rgb = PptxRGBColor(rgb[0], rgb[1], rgb[2])
                                new_run.font.bold = False
    out_io = BytesIO()
    prs.save(out_io)
    return out_io.getvalue()

# --- PDFチェック用関数 ---
def check_pdf(file, rules):
    results = []
    invisible_chars = r'[\s\u200B-\u200F\u202A-\u202E\u2060-\u206F\uFEFF\u00A0]+'
    
    with pdfplumber.open(file) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text(x_tolerance=2, y_tolerance=2)
            if not text: continue
            
            pure_text = re.sub(invisible_chars, '', text)
            parts = apply_rules_to_text(pure_text, rules, for_reporting=True)
            full_text = "".join([p[1] for p in parts])
            
            current_idx = 0
            for orig, curr, is_fixed, is_alnum in parts:
                if is_fixed:
                    start_idx = max(0, current_idx - 15)
                    end_idx = min(len(full_text), current_idx + len(curr) + 15)
                    context = full_text[start_idx:end_idx]
                    
                    reason = "英数字の半角化" if orig.translate(str.maketrans("ＡＢＣＤＥＦＧＨＩＪＫＬＭＮＯＰＱＲＳＴＵＶＷＸＹＺａｂｃｄｅｆｇｈｉｊｋｌｍｎｏｐｑｒｓｔｕｖｗｘｙｚ０１２３４５６７８９", "ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789")) == curr and orig != curr else "統一ルールの適用"
                    
                    results.append({
                        "ページ": i + 1,
                        "NGワード": orig,
                        "修正案": curr,
                        "修正理由": reason,
                        "周辺の文章": f"...{context}..."
                    })
                current_idx += len(curr)
                
    return results

# 4. メイン処理
uploaded_files = st.file_uploader("ファイルをアップロード (Word, Excel, PPT, PDF)", type=["docx", "xlsx", "pptx", "pdf"], accept_multiple_files=True)

if uploaded_files:
    for idx, file in enumerate(uploaded_files):
        ext = file.name.split('.')[-1].lower()
        if ext == "pdf":
            with st.spinner(f"PDF {file.name} を解析中..."):
                pdf_results = check_pdf(file, rules_dict)
                st.subheader(f"📑 {file.name} のチェック結果")
                if pdf_results:
                    st.warning(f"以下の {len(pdf_results)} 箇所で修正が推奨されます。")
                    st.table(pd.DataFrame(pdf_results))
                else:
                    st.success("統一ルールに反する箇所は見つかりませんでした！")
        else:
            with st.spinner(f"{file.name} を処理中..."):
                if ext == "docx": data = repair_docx(file, rules_dict, selected_rgb)
                elif ext == "xlsx": data = repair_xlsx(file, rules_dict, selected_rgb)
                elif ext == "pptx": data = repair_pptx(file, rules_dict, selected_rgb)
                
                st.download_button(label=f"📥 修正済みの {file.name} を保存", data=data, file_name=file.name, key=f"dl_btn_{idx}")
